"""Tests for fake file generators (Office, Image, XML)."""

import tempfile
from pathlib import Path

import pytest

from datadetector import FakeDataGenerator


@pytest.fixture
def generator():
    """Create FakeDataGenerator instance."""
    return FakeDataGenerator(seed=12345)


@pytest.fixture
def temp_dir():
    """Create temporary directory for test files."""
    with tempfile.TemporaryDirectory() as tmpdir:
        yield Path(tmpdir)


class TestOfficeFileGenerator:
    """Tests for Office file generation."""

    def test_create_word_file_basic(self, generator, temp_dir):
        """Test creating basic Word file."""
        pytest.importorskip("docx")
        from datadetector.fake_file_generators import OfficeFileGenerator

        office_gen = OfficeFileGenerator(generator)
        word_path = temp_dir / "test.docx"
        office_gen.create_word_file(word_path, paragraphs=5, include_pii=True)

        assert word_path.exists()
        assert word_path.stat().st_size > 0

    def test_create_word_file_no_pii(self, generator, temp_dir):
        """Test creating Word file without PII."""
        pytest.importorskip("docx")
        from datadetector.fake_file_generators import OfficeFileGenerator

        office_gen = OfficeFileGenerator(generator)
        word_path = temp_dir / "test_no_pii.docx"
        office_gen.create_word_file(word_path, paragraphs=3, include_pii=False)

        assert word_path.exists()

    def test_create_word_file_content(self, generator, temp_dir):
        """Test Word file content structure."""
        pytest.importorskip("docx")
        from datadetector.fake_file_generators import OfficeFileGenerator
        from docx import Document

        office_gen = OfficeFileGenerator(generator)
        word_path = temp_dir / "test_content.docx"
        office_gen.create_word_file(word_path, paragraphs=10, include_pii=True)

        # Read and verify
        doc = Document(word_path)
        assert len(doc.paragraphs) > 0
        # Check for title
        assert doc.paragraphs[0].text == "Fake Document with Data"

    def test_create_excel_file_basic(self, generator, temp_dir):
        """Test creating basic Excel file."""
        pytest.importorskip("openpyxl")
        from datadetector.fake_file_generators import OfficeFileGenerator

        office_gen = OfficeFileGenerator(generator)
        excel_path = temp_dir / "test.xlsx"
        office_gen.create_excel_file(excel_path, rows=50, include_pii=True)

        assert excel_path.exists()
        assert excel_path.stat().st_size > 0

    def test_create_excel_file_content(self, generator, temp_dir):
        """Test Excel file content structure."""
        pytest.importorskip("openpyxl")
        from datadetector.fake_file_generators import OfficeFileGenerator
        from openpyxl import load_workbook

        office_gen = OfficeFileGenerator(generator)
        excel_path = temp_dir / "test_content.xlsx"
        office_gen.create_excel_file(excel_path, rows=10, include_pii=True)

        # Read and verify
        wb = load_workbook(excel_path)
        ws = wb.active

        # Check headers
        headers = [cell.value for cell in ws[1]]
        assert "ID" in headers
        assert "Name" in headers
        assert "Email" in headers
        assert "SSN" in headers

        # Check row count (header + data rows)
        assert ws.max_row == 11  # 1 header + 10 data rows

    def test_create_excel_file_no_pii(self, generator, temp_dir):
        """Test Excel file without PII."""
        pytest.importorskip("openpyxl")
        from datadetector.fake_file_generators import OfficeFileGenerator
        from openpyxl import load_workbook

        office_gen = OfficeFileGenerator(generator)
        excel_path = temp_dir / "test_no_pii.xlsx"
        office_gen.create_excel_file(excel_path, rows=5, include_pii=False)

        wb = load_workbook(excel_path)
        ws = wb.active
        headers = [cell.value for cell in ws[1]]

        assert "SSN" not in headers
        assert "Credit Card" not in headers

    def test_create_excel_file_large(self, generator, temp_dir):
        """Test creating large Excel file."""
        pytest.importorskip("openpyxl")
        from datadetector.fake_file_generators import OfficeFileGenerator
        from openpyxl import load_workbook

        office_gen = OfficeFileGenerator(generator)
        excel_path = temp_dir / "test_large.xlsx"
        office_gen.create_excel_file(excel_path, rows=1000, include_pii=True)

        wb = load_workbook(excel_path)
        ws = wb.active
        assert ws.max_row == 1001  # 1 header + 1000 data rows

    def test_create_powerpoint_file_basic(self, generator, temp_dir):
        """Test creating basic PowerPoint file."""
        pytest.importorskip("pptx")
        from datadetector.fake_file_generators import OfficeFileGenerator

        office_gen = OfficeFileGenerator(generator)
        ppt_path = temp_dir / "test.pptx"
        office_gen.create_powerpoint_file(ppt_path, slides=5, include_pii=True)

        assert ppt_path.exists()
        assert ppt_path.stat().st_size > 0

    def test_create_powerpoint_file_content(self, generator, temp_dir):
        """Test PowerPoint file content structure."""
        pytest.importorskip("pptx")
        from datadetector.fake_file_generators import OfficeFileGenerator
        from pptx import Presentation

        office_gen = OfficeFileGenerator(generator)
        ppt_path = temp_dir / "test_content.pptx"
        office_gen.create_powerpoint_file(ppt_path, slides=3, include_pii=True)

        # Read and verify
        prs = Presentation(ppt_path)
        # Should have title slide + 3 content slides
        assert len(prs.slides) == 4

        # Check title slide
        title_slide = prs.slides[0]
        assert title_slide.shapes.title.text == "Fake Data Presentation"

    def test_create_powerpoint_file_no_pii(self, generator, temp_dir):
        """Test PowerPoint file without PII."""
        pytest.importorskip("pptx")
        from datadetector.fake_file_generators import OfficeFileGenerator

        office_gen = OfficeFileGenerator(generator)
        ppt_path = temp_dir / "test_no_pii.pptx"
        office_gen.create_powerpoint_file(ppt_path, slides=2, include_pii=False)

        assert ppt_path.exists()

    def test_office_generator_missing_dependency_word(self, generator):
        """Test error handling when python-docx is not available."""
        # This test would need to mock the import, skip for now
        pass

    def test_office_generator_missing_dependency_excel(self, generator):
        """Test error handling when openpyxl is not available."""
        # This test would need to mock the import, skip for now
        pass


class TestImageGenerator:
    """Tests for image generation."""

    def test_create_image_with_text_basic(self, generator, temp_dir):
        """Test creating basic image with text."""
        pytest.importorskip("PIL")
        from datadetector.fake_file_generators import ImageGenerator

        img_gen = ImageGenerator(generator)
        img_path = temp_dir / "test.png"
        img_gen.create_image_with_text(img_path, width=800, height=600, include_pii=True)

        assert img_path.exists()
        assert img_path.stat().st_size > 0

    def test_create_image_with_text_formats(self, generator, temp_dir):
        """Test creating images in different formats."""
        pytest.importorskip("PIL")
        from datadetector.fake_file_generators import ImageGenerator

        img_gen = ImageGenerator(generator)

        # Test PNG
        png_path = temp_dir / "test.png"
        img_gen.create_image_with_text(png_path, format="PNG")
        assert png_path.exists()

        # Test JPEG
        jpg_path = temp_dir / "test.jpg"
        img_gen.create_image_with_text(jpg_path, format="JPEG")
        assert jpg_path.exists()

    def test_create_image_with_text_sizes(self, generator, temp_dir):
        """Test creating images with different sizes."""
        pytest.importorskip("PIL")
        from datadetector.fake_file_generators import ImageGenerator
        from PIL import Image

        img_gen = ImageGenerator(generator)

        # Small image
        small_path = temp_dir / "small.png"
        img_gen.create_image_with_text(small_path, width=400, height=300)
        img = Image.open(small_path)
        assert img.size == (400, 300)

        # Large image
        large_path = temp_dir / "large.png"
        img_gen.create_image_with_text(large_path, width=1920, height=1080)
        img = Image.open(large_path)
        assert img.size == (1920, 1080)

    def test_create_image_with_text_no_pii(self, generator, temp_dir):
        """Test creating image without PII."""
        pytest.importorskip("PIL")
        from datadetector.fake_file_generators import ImageGenerator

        img_gen = ImageGenerator(generator)
        img_path = temp_dir / "test_no_pii.png"
        img_gen.create_image_with_text(img_path, include_pii=False)

        assert img_path.exists()

    def test_create_screenshot_like_image_basic(self, generator, temp_dir):
        """Test creating screenshot-like image."""
        pytest.importorskip("PIL")
        from datadetector.fake_file_generators import ImageGenerator

        img_gen = ImageGenerator(generator)
        img_path = temp_dir / "screenshot.png"
        img_gen.create_screenshot_like_image(img_path, include_pii=True)

        assert img_path.exists()
        assert img_path.stat().st_size > 0

    def test_create_screenshot_like_image_content(self, generator, temp_dir):
        """Test screenshot-like image size and format."""
        pytest.importorskip("PIL")
        from datadetector.fake_file_generators import ImageGenerator
        from PIL import Image

        img_gen = ImageGenerator(generator)
        img_path = temp_dir / "screenshot_content.png"
        img_gen.create_screenshot_like_image(img_path, include_pii=True)

        img = Image.open(img_path)
        assert img.size == (1200, 800)
        assert img.mode == "RGB"

    def test_create_screenshot_like_image_no_pii(self, generator, temp_dir):
        """Test screenshot-like image without PII."""
        pytest.importorskip("PIL")
        from datadetector.fake_file_generators import ImageGenerator

        img_gen = ImageGenerator(generator)
        img_path = temp_dir / "screenshot_no_pii.png"
        img_gen.create_screenshot_like_image(img_path, include_pii=False)

        assert img_path.exists()

    def test_image_generator_missing_dependency(self, generator):
        """Test error handling when PIL is not available."""
        # Would need to mock PIL import to test this
        pass


class TestXMLGenerator:
    """Tests for XML generation."""

    def test_create_xml_file_basic(self, generator, temp_dir):
        """Test creating basic XML file."""
        from datadetector.fake_file_generators import XMLGenerator

        xml_gen = XMLGenerator(generator)
        xml_path = temp_dir / "test.xml"
        xml_gen.create_xml_file(xml_path, records=10, include_pii=True)

        assert xml_path.exists()
        assert xml_path.stat().st_size > 0

    def test_create_xml_file_structure(self, generator, temp_dir):
        """Test XML file structure."""
        import xml.etree.ElementTree as ET

        from datadetector.fake_file_generators import XMLGenerator

        xml_gen = XMLGenerator(generator)
        xml_path = temp_dir / "test_structure.xml"
        xml_gen.create_xml_file(xml_path, records=5, include_pii=True)

        # Parse and verify
        tree = ET.parse(xml_path)
        root = tree.getroot()

        assert root.tag == "users"
        assert len(root) == 5  # 5 user records

        # Check first user
        user = root[0]
        assert user.tag == "user"
        assert user.find("name") is not None
        assert user.find("email") is not None
        assert user.find("ssn") is not None

    def test_create_xml_file_no_pii(self, generator, temp_dir):
        """Test XML file without PII."""
        import xml.etree.ElementTree as ET

        from datadetector.fake_file_generators import XMLGenerator

        xml_gen = XMLGenerator(generator)
        xml_path = temp_dir / "test_no_pii.xml"
        xml_gen.create_xml_file(xml_path, records=3, include_pii=False)

        tree = ET.parse(xml_path)
        root = tree.getroot()
        user = root[0]

        assert user.find("ssn") is None
        assert user.find("credit_card") is None
        assert user.find("name") is not None
        assert user.find("email") is not None

    def test_create_xml_file_attributes(self, generator, temp_dir):
        """Test XML file attributes."""
        import xml.etree.ElementTree as ET

        from datadetector.fake_file_generators import XMLGenerator

        xml_gen = XMLGenerator(generator)
        xml_path = temp_dir / "test_attrs.xml"
        xml_gen.create_xml_file(xml_path, records=5, include_pii=True)

        tree = ET.parse(xml_path)
        root = tree.getroot()

        # Check user IDs
        for i, user in enumerate(root, 1):
            assert user.get("id") == str(i)

    def test_create_xml_file_large(self, generator, temp_dir):
        """Test creating large XML file."""
        import xml.etree.ElementTree as ET

        from datadetector.fake_file_generators import XMLGenerator

        xml_gen = XMLGenerator(generator)
        xml_path = temp_dir / "test_large.xml"
        xml_gen.create_xml_file(xml_path, records=1000, include_pii=True)

        tree = ET.parse(xml_path)
        root = tree.getroot()
        assert len(root) == 1000


class TestPDFGenerator:
    """Tests for PDF generation."""

    def test_create_pdf_file_basic(self, generator, temp_dir):
        """Test creating basic PDF file."""
        pytest.importorskip("reportlab")
        from datadetector.fake_file_generators import PDFGenerator

        pdf_gen = PDFGenerator(generator)
        pdf_path = temp_dir / "test.pdf"
        pdf_gen.create_pdf_file(pdf_path, pages=5, include_pii=True)

        assert pdf_path.exists()
        assert pdf_path.stat().st_size > 0

    def test_create_pdf_file_no_pii(self, generator, temp_dir):
        """Test creating PDF without PII."""
        pytest.importorskip("reportlab")
        from datadetector.fake_file_generators import PDFGenerator

        pdf_gen = PDFGenerator(generator)
        pdf_path = temp_dir / "test_no_pii.pdf"
        pdf_gen.create_pdf_file(pdf_path, pages=3, include_pii=False)

        assert pdf_path.exists()

    def test_create_pdf_file_pages(self, generator, temp_dir):
        """Test creating PDF with different page counts."""
        pytest.importorskip("reportlab")
        from datadetector.fake_file_generators import PDFGenerator

        pdf_gen = PDFGenerator(generator)

        # Single page
        single_path = temp_dir / "single.pdf"
        pdf_gen.create_pdf_file(single_path, pages=1, include_pii=True)
        assert single_path.exists()

        # Multiple pages
        multi_path = temp_dir / "multi.pdf"
        pdf_gen.create_pdf_file(multi_path, pages=10, include_pii=True)
        assert multi_path.exists()
        # More pages should generally result in larger file
        assert multi_path.stat().st_size > single_path.stat().st_size

    def test_create_pdf_invoice_basic(self, generator, temp_dir):
        """Test creating PDF invoice."""
        pytest.importorskip("reportlab")
        from datadetector.fake_file_generators import PDFGenerator

        pdf_gen = PDFGenerator(generator)
        invoice_path = temp_dir / "invoice.pdf"
        pdf_gen.create_pdf_invoice(invoice_path, include_pii=True)

        assert invoice_path.exists()
        assert invoice_path.stat().st_size > 0

    def test_create_pdf_invoice_no_pii(self, generator, temp_dir):
        """Test creating PDF invoice without PII."""
        pytest.importorskip("reportlab")
        from datadetector.fake_file_generators import PDFGenerator

        pdf_gen = PDFGenerator(generator)
        invoice_path = temp_dir / "invoice_no_pii.pdf"
        pdf_gen.create_pdf_invoice(invoice_path, include_pii=False)

        assert invoice_path.exists()

    def test_pdf_generator_missing_dependency(self, generator):
        """Test error handling when reportlab is not available."""
        # Would need to mock reportlab import to test this properly
        # For now, just verify the generator can be instantiated
        from datadetector.fake_file_generators import PDFGenerator

        pdf_gen = PDFGenerator(generator)
        assert pdf_gen.faker is generator.faker


class TestGeneratorsIntegration:
    """Tests for integration between different generators."""

    def test_all_generators_use_same_faker(self, generator, temp_dir):
        """Test that all generators share the same Faker instance."""
        pytest.importorskip("docx")
        pytest.importorskip("openpyxl")
        pytest.importorskip("pptx")
        pytest.importorskip("PIL")

        from datadetector.fake_file_generators import (
            ImageGenerator,
            OfficeFileGenerator,
            XMLGenerator,
        )

        office_gen = OfficeFileGenerator(generator)
        img_gen = ImageGenerator(generator)
        xml_gen = XMLGenerator(generator)

        # All should reference the same faker instance
        assert office_gen.faker is generator.faker
        assert img_gen.faker is generator.faker
        assert xml_gen.faker is generator.faker

    def test_generators_with_detector(self, generator, temp_dir):
        """Test that generated files can be processed by detector."""
        pytest.importorskip("openpyxl")
        from datadetector import Engine, load_registry
        from datadetector.fake_file_generators import OfficeFileGenerator, XMLGenerator
        from openpyxl import load_workbook

        registry = load_registry()
        engine = Engine(registry)

        # Generate Excel file
        office_gen = OfficeFileGenerator(generator)
        excel_path = temp_dir / "detect_test.xlsx"
        office_gen.create_excel_file(excel_path, rows=10, include_pii=True)

        # Read Excel and detect PII
        wb = load_workbook(excel_path)
        ws = wb.active

        # Get email from first data row
        email_col = None
        for idx, cell in enumerate(ws[1], 1):
            if cell.value == "Email":
                email_col = idx
                break

        if email_col:
            email = ws.cell(row=2, column=email_col).value
            result = engine.validate(email, "comm/email_01")
            assert result.is_valid

        # Generate XML file
        xml_gen = XMLGenerator(generator)
        xml_path = temp_dir / "detect_test.xml"
        xml_gen.create_xml_file(xml_path, records=5, include_pii=True)

        # Read XML and detect
        with open(xml_path, "r", encoding="utf-8") as f:
            xml_content = f.read()

        result = engine.find(xml_content)
        assert result.has_matches
