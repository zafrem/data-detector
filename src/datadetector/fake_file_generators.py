"""File format generators for office documents and images.

This module extends the fake data generator with support for:
- Microsoft Office files (Word, Excel, PowerPoint)
- Images with embedded text
- PDF files
- XML files
"""

import io
import logging
import random
from pathlib import Path
from typing import List, Optional, Union

try:
    from PIL import Image, ImageDraw, ImageFont
except ImportError:
    Image = ImageDraw = ImageFont = None

logger = logging.getLogger(__name__)


class OfficeFileGenerator:
    """Generate Office files (Word, Excel, PowerPoint) with fake data."""

    def __init__(self, faker_generator):
        """
        Initialize with a FakeDataGenerator instance.

        Args:
            faker_generator: FakeDataGenerator instance
        """
        self.gen = faker_generator
        self.faker = faker_generator.faker

    def create_word_file(
        self,
        output_path: Union[str, Path],
        paragraphs: int = 10,
        include_pii: bool = True,
    ) -> None:
        """
        Create a Word document (.docx) with fake data.

        Args:
            output_path: Output file path
            paragraphs: Number of paragraphs
            include_pii: Whether to include PII

        Raises:
            ImportError: If python-docx is not installed
        """
        try:
            from docx import Document
        except ImportError:
            raise ImportError(
                "python-docx is required for Word file generation. "
                "Install it with: pip install python-docx"
            )

        output_path = Path(output_path)
        doc = Document()

        # Add title
        doc.add_heading('Fake Document with Data', 0)

        # Add paragraphs
        for i in range(paragraphs):
            text = self.faker.paragraph(nb_sentences=5)

            if include_pii and random.random() > 0.6:
                # Add a table with PII data
                doc.add_heading(f'Section {i+1}: Contact Information', level=2)

                table = doc.add_table(rows=5, cols=2)
                table.style = 'Light Grid Accent 1'

                cells = [
                    ('Name', self.faker.name()),
                    ('Email', self.faker.email()),
                    ('Phone', self.faker.phone_number()),
                    ('SSN', self.faker.ssn() if include_pii else 'XXX-XX-XXXX'),
                    ('Address', self.faker.address()),
                ]

                for idx, (label, value) in enumerate(cells):
                    table.rows[idx].cells[0].text = label
                    table.rows[idx].cells[1].text = value
            else:
                doc.add_paragraph(text)

        doc.save(output_path)
        logger.info(f"Created Word file: {output_path}")

    def create_excel_file(
        self,
        output_path: Union[str, Path],
        rows: int = 100,
        include_pii: bool = True,
    ) -> None:
        """
        Create an Excel file (.xlsx) with fake data.

        Args:
            output_path: Output file path
            rows: Number of data rows
            include_pii: Whether to include PII

        Raises:
            ImportError: If openpyxl is not installed
        """
        try:
            from openpyxl import Workbook
            from openpyxl.styles import Font, PatternFill
        except ImportError:
            raise ImportError(
                "openpyxl is required for Excel file generation. "
                "Install it with: pip install openpyxl"
            )

        output_path = Path(output_path)
        wb = Workbook()
        ws = wb.active
        ws.title = "Users"

        # Headers
        headers = ['ID', 'Name', 'Email', 'Phone', 'City', 'Country']
        if include_pii:
            headers.extend(['SSN', 'Credit Card', 'IP Address'])

        # Style headers
        header_fill = PatternFill(start_color="366092", end_color="366092", fill_type="solid")
        header_font = Font(bold=True, color="FFFFFF")

        for col_idx, header in enumerate(headers, 1):
            cell = ws.cell(row=1, column=col_idx)
            cell.value = header
            cell.fill = header_fill
            cell.font = header_font

        # Data rows
        for row_idx in range(2, rows + 2):
            ws.cell(row=row_idx, column=1, value=row_idx - 1)
            ws.cell(row=row_idx, column=2, value=self.faker.name())
            ws.cell(row=row_idx, column=3, value=self.faker.email())
            ws.cell(row=row_idx, column=4, value=self.faker.phone_number())
            ws.cell(row=row_idx, column=5, value=self.faker.city())
            ws.cell(row=row_idx, column=6, value=self.faker.country())

            if include_pii:
                ws.cell(row=row_idx, column=7, value=self.faker.ssn())
                ws.cell(row=row_idx, column=8, value=self.faker.credit_card_number())
                ws.cell(row=row_idx, column=9, value=self.faker.ipv4())

        # Adjust column widths
        for column in ws.columns:
            max_length = 0
            column = list(column)
            for cell in column:
                try:
                    if len(str(cell.value)) > max_length:
                        max_length = len(cell.value)
                except:
                    pass
            adjusted_width = min(max_length + 2, 50)
            ws.column_dimensions[column[0].column_letter].width = adjusted_width

        wb.save(output_path)
        logger.info(f"Created Excel file: {output_path} ({rows} rows)")

    def create_powerpoint_file(
        self,
        output_path: Union[str, Path],
        slides: int = 5,
        include_pii: bool = True,
    ) -> None:
        """
        Create a PowerPoint file (.pptx) with fake data.

        Args:
            output_path: Output file path
            slides: Number of slides
            include_pii: Whether to include PII

        Raises:
            ImportError: If python-pptx is not installed
        """
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            raise ImportError(
                "python-pptx is required for PowerPoint generation. "
                "Install it with: pip install python-pptx"
            )

        output_path = Path(output_path)
        prs = Presentation()

        # Title slide
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]

        title.text = "Fake Data Presentation"
        subtitle.text = f"Generated with {slides} slides"

        # Content slides
        for i in range(slides):
            bullet_slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(bullet_slide_layout)

            title = slide.shapes.title
            body = slide.placeholders[1]

            title.text = f"Slide {i+1}: {self.faker.catch_phrase()}"

            tf = body.text_frame
            tf.text = self.faker.paragraph()

            if include_pii and random.random() > 0.5:
                # Add PII data
                p = tf.add_paragraph()
                p.text = f"Contact: {self.faker.email()}"
                p.level = 1

                p = tf.add_paragraph()
                p.text = f"Phone: {self.faker.phone_number()}"
                p.level = 1

        prs.save(output_path)
        logger.info(f"Created PowerPoint file: {output_path} ({slides} slides)")


class ImageGenerator:
    """Generate images with embedded text/PII."""

    def __init__(self, faker_generator):
        """
        Initialize with a FakeDataGenerator instance.

        Args:
            faker_generator: FakeDataGenerator instance
        """
        if Image is None or ImageDraw is None:
            raise ImportError(
                "PIL/Pillow is required for image generation. "
                "Install it with: pip install Pillow"
            )

        self.gen = faker_generator
        self.faker = faker_generator.faker

    def create_image_with_text(
        self,
        output_path: Union[str, Path],
        width: int = 800,
        height: int = 600,
        include_pii: bool = True,
        format: str = "PNG",
    ) -> None:
        """
        Create an image with embedded text.

        Args:
            output_path: Output file path
            width: Image width in pixels
            height: Image height in pixels
            include_pii: Whether to include PII in text
            format: Image format (PNG, JPEG, etc.)
        """
        output_path = Path(output_path)

        # Create image
        img = Image.new('RGB', (width, height), color=(255, 255, 255))
        draw = ImageDraw.Draw(img)

        # Try to use a nice font, fall back to default
        try:
            font_large = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 32)
            font_small = ImageFont.truetype("/System/Library/Fonts/Helvetica.ttc", 20)
        except:
            font_large = ImageFont.load_default()
            font_small = ImageFont.load_default()

        # Add title
        title = "Fake Document"
        draw.text((50, 50), title, fill=(0, 0, 0), font=font_large)

        # Add content
        y_position = 120
        lines = []

        if include_pii:
            lines = [
                f"Name: {self.faker.name()}",
                f"Email: {self.faker.email()}",
                f"Phone: {self.faker.phone_number()}",
                f"Address: {self.faker.address()}",
                f"SSN: {self.faker.ssn()}",
                f"Credit Card: {self.faker.credit_card_number()}",
            ]
        else:
            lines = [
                f"Company: {self.faker.company()}",
                f"Job: {self.faker.job()}",
                f"Date: {self.faker.date()}",
            ]

        for line in lines:
            draw.text((50, y_position), line, fill=(0, 0, 0), font=font_small)
            y_position += 40

        img.save(output_path, format=format)
        logger.info(f"Created image: {output_path} ({width}x{height})")

    def create_screenshot_like_image(
        self,
        output_path: Union[str, Path],
        include_pii: bool = True,
    ) -> None:
        """
        Create an image that looks like a screenshot with data.

        Args:
            output_path: Output file path
            include_pii: Whether to include PII
        """
        output_path = Path(output_path)

        width, height = 1200, 800
        img = Image.new('RGB', (width, height), color=(240, 240, 240))
        draw = ImageDraw.Draw(img)

        # Try to load font
        try:
            font = ImageFont.truetype("/System/Library/Fonts/Menlo.ttc", 14)
        except:
            font = ImageFont.load_default()

        # Draw fake terminal/editor look
        draw.rectangle([(50, 50), (width-50, height-50)], fill=(30, 30, 30))

        # Add fake code/config with PII
        y_pos = 80
        lines = [
            "# Configuration File",
            "",
            f"USERNAME={self.faker.user_name()}",
            f"EMAIL={self.faker.email()}",
        ]

        if include_pii:
            lines.extend([
                f"API_KEY={self.gen.from_pattern('comm/aws_access_key_01')}",
                f"DB_PASSWORD={self.faker.password()}",
                f"SSN={self.faker.ssn()}",
            ])

        for line in lines:
            draw.text((70, y_pos), line, fill=(0, 255, 0), font=font)
            y_pos += 25

        img.save(output_path, format="PNG")
        logger.info(f"Created screenshot-like image: {output_path}")


class XMLGenerator:
    """Generate XML files with fake data."""

    def __init__(self, faker_generator):
        """
        Initialize with a FakeDataGenerator instance.

        Args:
            faker_generator: FakeDataGenerator instance
        """
        self.gen = faker_generator
        self.faker = faker_generator.faker

    def create_xml_file(
        self,
        output_path: Union[str, Path],
        records: int = 50,
        include_pii: bool = True,
    ) -> None:
        """
        Create an XML file with fake data.

        Args:
            output_path: Output file path
            records: Number of records
            include_pii: Whether to include PII
        """
        try:
            import xml.etree.ElementTree as ET
        except ImportError:
            raise ImportError("xml.etree.ElementTree is required")

        output_path = Path(output_path)

        root = ET.Element("users")

        for i in range(records):
            user = ET.SubElement(root, "user", id=str(i+1))

            ET.SubElement(user, "name").text = self.faker.name()
            ET.SubElement(user, "email").text = self.faker.email()
            ET.SubElement(user, "phone").text = self.faker.phone_number()

            if include_pii:
                ET.SubElement(user, "ssn").text = self.faker.ssn()
                ET.SubElement(user, "credit_card").text = self.faker.credit_card_number()

        tree = ET.ElementTree(root)
        tree.write(output_path, encoding='utf-8', xml_declaration=True)

        logger.info(f"Created XML file: {output_path} ({records} records)")


class PDFGenerator:
    """Generate PDF files with fake data."""

    def __init__(self, faker_generator):
        """
        Initialize with a FakeDataGenerator instance.

        Args:
            faker_generator: FakeDataGenerator instance
        """
        self.gen = faker_generator
        self.faker = faker_generator.faker

    def create_pdf_file(
        self,
        output_path: Union[str, Path],
        pages: int = 5,
        include_pii: bool = True,
    ) -> None:
        """
        Create a PDF file with fake data.

        Args:
            output_path: Output file path
            pages: Number of pages to generate
            include_pii: Whether to include PII

        Raises:
            ImportError: If reportlab is not installed
        """
        try:
            from reportlab.lib import colors
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import ParagraphStyle, getSampleStyleSheet
            from reportlab.lib.units import inch
            from reportlab.platypus import (
                PageBreak,
                Paragraph,
                SimpleDocTemplate,
                Spacer,
                Table,
                TableStyle,
            )
        except ImportError:
            raise ImportError(
                "reportlab is required for PDF generation. "
                "Install it with: pip install reportlab"
            )

        output_path = Path(output_path)

        # Create PDF document
        doc = SimpleDocTemplate(str(output_path), pagesize=letter)
        story = []
        styles = getSampleStyleSheet()

        # Custom styles
        title_style = ParagraphStyle(
            'CustomTitle',
            parent=styles['Heading1'],
            fontSize=24,
            textColor=colors.HexColor('#2C3E50'),
            spaceAfter=30,
        )

        heading_style = ParagraphStyle(
            'CustomHeading',
            parent=styles['Heading2'],
            fontSize=16,
            textColor=colors.HexColor('#34495E'),
            spaceAfter=12,
        )

        # Title page
        story.append(Paragraph("Fake Document with Data", title_style))
        story.append(Spacer(1, 0.2 * inch))
        story.append(Paragraph(f"Generated: {self.faker.date()}", styles['Normal']))
        story.append(Spacer(1, 0.5 * inch))

        # Generate pages
        for page_num in range(pages):
            if page_num > 0:
                story.append(PageBreak())

            # Page heading
            story.append(Paragraph(f"Section {page_num + 1}: {self.faker.catch_phrase()}", heading_style))
            story.append(Spacer(1, 0.2 * inch))

            # Add paragraphs
            for _ in range(random.randint(2, 4)):
                paragraph = self.faker.paragraph(nb_sentences=random.randint(3, 6))
                story.append(Paragraph(paragraph, styles['Normal']))
                story.append(Spacer(1, 0.1 * inch))

            # Add PII table on some pages
            if include_pii and random.random() > 0.4:
                story.append(Spacer(1, 0.2 * inch))
                story.append(Paragraph("Contact Information", heading_style))
                story.append(Spacer(1, 0.1 * inch))

                # Create table data
                table_data = [
                    ['Field', 'Value'],
                    ['Name', self.faker.name()],
                    ['Email', self.faker.email()],
                    ['Phone', self.faker.phone_number()],
                    ['Address', self.faker.address().replace('\n', ', ')],
                ]

                if include_pii:
                    table_data.extend([
                        ['SSN', self.faker.ssn()],
                        ['Credit Card', self.faker.credit_card_number()],
                        ['IP Address', self.faker.ipv4()],
                    ])

                # Create table
                table = Table(table_data, colWidths=[2 * inch, 4 * inch])
                table.setStyle(TableStyle([
                    ('BACKGROUND', (0, 0), (-1, 0), colors.HexColor('#3498DB')),
                    ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
                    ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
                    ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
                    ('FONTSIZE', (0, 0), (-1, 0), 12),
                    ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
                    ('BACKGROUND', (0, 1), (-1, -1), colors.beige),
                    ('GRID', (0, 0), (-1, -1), 1, colors.black),
                    ('FONTNAME', (0, 1), (-1, -1), 'Helvetica'),
                    ('FONTSIZE', (0, 1), (-1, -1), 10),
                    ('ROWBACKGROUNDS', (0, 1), (-1, -1), [colors.white, colors.lightgrey]),
                ]))

                story.append(table)
                story.append(Spacer(1, 0.3 * inch))

            # Add some bullet points
            if random.random() > 0.5:
                story.append(Spacer(1, 0.2 * inch))
                story.append(Paragraph("Key Points:", heading_style))
                story.append(Spacer(1, 0.1 * inch))

                for _ in range(random.randint(3, 5)):
                    bullet = f"• {self.faker.sentence()}"
                    story.append(Paragraph(bullet, styles['Normal']))
                    story.append(Spacer(1, 0.05 * inch))

        # Build PDF
        doc.build(story)
        logger.info(f"Created PDF file: {output_path} ({pages} pages)")

    def create_pdf_invoice(
        self,
        output_path: Union[str, Path],
        include_pii: bool = True,
    ) -> None:
        """
        Create a PDF invoice with fake data.

        Args:
            output_path: Output file path
            include_pii: Whether to include PII

        Raises:
            ImportError: If reportlab is not installed
        """
        try:
            from reportlab.lib import colors
            from reportlab.lib.pagesizes import letter
            from reportlab.lib.styles import getSampleStyleSheet
            from reportlab.lib.units import inch
            from reportlab.platypus import (
                Paragraph,
                SimpleDocTemplate,
                Spacer,
                Table,
                TableStyle,
            )
        except ImportError:
            raise ImportError(
                "reportlab is required for PDF generation. "
                "Install it with: pip install reportlab"
            )

        output_path = Path(output_path)
        doc = SimpleDocTemplate(str(output_path), pagesize=letter)
        story = []
        styles = getSampleStyleSheet()

        # Invoice header
        story.append(Paragraph("INVOICE", styles['Title']))
        story.append(Spacer(1, 0.3 * inch))

        # Invoice info
        invoice_num = f"INV-{random.randint(10000, 99999)}"
        story.append(Paragraph(f"Invoice Number: {invoice_num}", styles['Normal']))
        story.append(Paragraph(f"Date: {self.faker.date()}", styles['Normal']))
        story.append(Spacer(1, 0.3 * inch))

        # Billing information
        story.append(Paragraph("Bill To:", styles['Heading2']))
        story.append(Paragraph(self.faker.name(), styles['Normal']))
        if include_pii:
            story.append(Paragraph(self.faker.email(), styles['Normal']))
            story.append(Paragraph(self.faker.phone_number(), styles['Normal']))
        story.append(Paragraph(self.faker.address().replace('\n', '<br/>'), styles['Normal']))
        story.append(Spacer(1, 0.3 * inch))

        # Items table
        story.append(Paragraph("Items:", styles['Heading2']))
        story.append(Spacer(1, 0.1 * inch))

        items_data = [['Item', 'Quantity', 'Price', 'Total']]
        total = 0
        for _ in range(random.randint(3, 8)):
            item = self.faker.catch_phrase()
            qty = random.randint(1, 10)
            price = random.uniform(10, 500)
            item_total = qty * price
            total += item_total
            items_data.append([
                item,
                str(qty),
                f"${price:.2f}",
                f"${item_total:.2f}"
            ])

        items_data.append(['', '', 'Subtotal:', f"${total:.2f}"])
        tax = total * 0.08
        items_data.append(['', '', 'Tax (8%):', f"${tax:.2f}"])
        items_data.append(['', '', 'Total:', f"${total + tax:.2f}"])

        items_table = Table(items_data, colWidths=[3 * inch, 1 * inch, 1.25 * inch, 1.25 * inch])
        items_table.setStyle(TableStyle([
            ('BACKGROUND', (0, 0), (-1, 0), colors.grey),
            ('TEXTCOLOR', (0, 0), (-1, 0), colors.whitesmoke),
            ('ALIGN', (0, 0), (-1, -1), 'LEFT'),
            ('ALIGN', (1, 1), (-1, -1), 'RIGHT'),
            ('FONTNAME', (0, 0), (-1, 0), 'Helvetica-Bold'),
            ('FONTSIZE', (0, 0), (-1, 0), 12),
            ('BOTTOMPADDING', (0, 0), (-1, 0), 12),
            ('BACKGROUND', (0, 1), (-1, -4), colors.beige),
            ('GRID', (0, 0), (-1, -4), 1, colors.black),
            ('LINEABOVE', (2, -3), (-1, -3), 1, colors.black),
            ('LINEABOVE', (2, -1), (-1, -1), 2, colors.black),
            ('FONTNAME', (2, -1), (-1, -1), 'Helvetica-Bold'),
        ]))

        story.append(items_table)
        story.append(Spacer(1, 0.5 * inch))

        # Payment info
        if include_pii:
            story.append(Paragraph("Payment Information:", styles['Heading2']))
            story.append(Paragraph(f"Account: {self.faker.credit_card_number()}", styles['Normal']))

        doc.build(story)
        logger.info(f"Created PDF invoice: {output_path}")
